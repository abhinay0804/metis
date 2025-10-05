"""
PowerPoint content extractor using python-pptx.
"""

from pptx import Presentation
from typing import Dict, Any, List
import json


class PPTExtractor:
    """Extract text, tables, and metadata from PowerPoint files."""
    
    def extract(self, file_path: str) -> Dict[str, Any]:
        """
        Extract content from PowerPoint file.
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            Dictionary containing extracted content
        """
        content = {
            "slides": [],
            "metadata": {},
            "slide_count": 0
        }
        
        try:
            prs = Presentation(file_path)
            content["slide_count"] = len(prs.slides)
            content["metadata"] = {
                "title": prs.core_properties.title,
                "author": prs.core_properties.author,
                "subject": prs.core_properties.subject,
                "created": str(prs.core_properties.created) if prs.core_properties.created else None,
                "modified": str(prs.core_properties.modified) if prs.core_properties.modified else None
            }
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = {
                    "slide_number": slide_num,
                    "title": "",
                    "text_content": [],
                    "tables": [],
                    "shapes": []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if shape.shape_type == 1:  # Title shape
                            slide_content["title"] = shape.text.strip()
                        else:
                            slide_content["text_content"].append({
                                "text": shape.text.strip(),
                                "shape_type": str(shape.shape_type)
                            })
                    
                    # Extract tables
                    if shape.has_table:
                        table_data = []
                        for row in shape.table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        
                        slide_content["tables"].append({
                            "data": table_data
                        })
                    
                    # Extract shape information
                    slide_content["shapes"].append({
                        "shape_type": str(shape.shape_type),
                        "has_text": hasattr(shape, "text") and bool(shape.text.strip()),
                        "has_table": shape.has_table
                    })
                
                content["slides"].append(slide_content)
        
        except Exception as e:
            raise Exception(f"Error extracting PowerPoint content: {str(e)}")
        
        return content
